import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def process_pdf(pdf_path, img_folder):
    doc = fitz.open(pdf_path)
    latex_content = []
    for i, page in enumerate(doc, start=1):
        latex_content.append(f"\\section*{{Pagina {i}}}")
        text = page.get_text("text")
        if text.strip():
            latex_content.append(text.replace("&", "\\&").replace("%", "\\%").replace("$", "\\$"))
        
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            img_filename = f"pdf_img_{i}_{img_index}.{base_image['ext']}"
            img_path = os.path.join(img_folder, img_filename)
            with open(img_path, "wb") as f:
                f.write(base_image["image"])
            latex_content.append(f"\\begin{{figure}}[h!]\\centering\\includegraphics[width=0.7\\textwidth]{{{img_path}}}\\end{{figure}}")
        latex_content.append("\\newpage")
    return "\n".join(latex_content)

def process_pptx(pptx_path, img_folder):
    prs = Presentation(pptx_path)
    latex_content = []
    img_counter = 0
    for i, slide in enumerate(prs.slides, start=1):
        latex_content.append(f"\\section*{{Slide {i}}}")
        shapes = sorted(slide.shapes, key=lambda s: (s.top, s.left))
        for shape in shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    latex_content.append(text.replace("&", "\\&").replace("%", "\\%").replace("$", "\\$"))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_counter += 1
                img_name = f"pptx_img_{i}_{img_counter}.{shape.image.ext}"
                img_path = os.path.join(img_folder, img_name)
                with open(img_path, "wb") as f:
                    f.write(shape.image.blob)
                latex_content.append(f"\\begin{{figure}}[h!]\\centering\\includegraphics[width=0.7\\textwidth]{{{img_path}}}\\end{{figure}}")
        latex_content.append("\\newpage")
    return "\n".join(latex_content)

def main():
    img_folder = "images"
    if not os.path.exists(img_folder): os.makedirs(img_folder)
    
    # 1. Trova i file disponibili
    files = [f for f in os.listdir('.') if f.endswith(('.pdf', '.pptx'))]
    if not files:
        print("❌ Nessun file trovato!")
        return

    print("\n--- 🎓 SlideToLaTeX Transcriber ---")
    print("File trovati nella cartella:")
    for i, f in enumerate(files, 1):
        print(f"{i}. {f}")
    
    # 2. Scelta del file da convertire
    scelta = int(input("\nInserisci il numero del file da convertire: ")) - 1
    target = files[scelta]
    
    # 3. Scelta del nome di output
    out_name = input("Come vuoi chiamare il file LaTeX? (default: output.tex): ").strip()
    if not out_name: out_name = "output.tex"
    if not out_name.endswith(".tex"): out_name += ".tex"

    # 4. Elaborazione
    ext = os.path.splitext(target)[1].lower()
    header = "\\documentclass{article}\\usepackage[utf8]{inputenc}\\usepackage{graphicx}\\begin{document}\n"
    footer = "\n\\end{document}"
    
    print(f"🚀 Elaborazione di {target}...")
    body = process_pptx(target, img_folder) if ext == ".pptx" else process_pdf(target, img_folder)

    with open(out_name, "w", encoding="utf-8") as f:
        f.write(header + body + footer)
    
    print(f"✅ Fatto! Creato '{out_name}' e salvate immagini in '{img_folder}'.")

if __name__ == "__main__":
    main()